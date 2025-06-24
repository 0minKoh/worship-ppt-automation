# utils/update_pptx.py

from pptx import Presentation
# PresentationType 대신 Presentation을 직접 사용하거나, 타입을 더 명시적으로 지정
from pptx.util import Inches
# from pptx.enum.shapes import MSO_SHAPE, MSO_AUTO_SIZE # MSO_AUTO_SIZE 임포트
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN # PP_ALIGN 임포트

# Type hint
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide as SlideType

def load_template(template_path: str) -> Presentation:
    """
    지정된 경로에서 PPTX 템플릿을 로드합니다.
    """
    try:
        prs = Presentation(template_path)
        print(f"Loaded template from: {template_path}")
        return prs
    except Exception as e:
        print(f"Error loading PPT template {template_path}: {e}")
        raise


def edit_text_field(
    prs: PresentationType,
    slide_index: int,
    new_text: str,
    is_title: bool = False,
    shape_name: str = None, # shape_name을 통한 접근 (PPTX 내부 이름)
    ph_index: int = None,    # placeholder_format.idx를 통한 접근 (PPTX 내부 인덱스)
    align_center: bool = True # 텍스트 중앙 정렬 여부
) -> PresentationType:
    """
    주어진 슬라이드의 텍스트 필드를 수정합니다.
    `is_title=True`이면 슬라이드의 제목 Placeholder를 찾고,
    `shape_name`이 주어지면 해당 이름의 도형을,
    `ph_index`가 주어지면 해당 인덱스의 Placeholder를 찾습니다.
    텍스트 프레임을 찾지 못하면 경고를 출력합니다.
    """
    if slide_index >= len(prs.slides):
        print(f"Warning: Slide index {slide_index} out of bounds for editing.")
        return prs

    slide = prs.slides[slide_index]
    found_text_frame = False
    text_frame = None

    if is_title:
        if slide.shapes.title:
            text_frame = slide.shapes.title.text_frame
        else:
            print(f"Warning: No title placeholder found on slide {slide_index} for title editing.")
    elif shape_name:
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                text_frame = shape.text_frame
                break
        if not text_frame:
            print(f"Warning: No shape named '{shape_name}' with text frame found on slide {slide_index}.")
    elif ph_index is not None:
        for shape in slide.shapes.placeholders:
            if shape.has_text_frame and shape.placeholder_format.idx == ph_index:
                text_frame = shape.text_frame
                break
        if not text_frame:
            print(f"Warning: No placeholder with index {ph_index} and text frame found on slide {slide_index}.")
    else: # Fallback: is_title도 아니고, shape_name, ph_index도 없으면, 첫 번째 텍스트 프레임 찾기
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                break
        if not text_frame:
            print(f"Warning: No generic text frame found on slide {slide_index} for editing.")
    
    if text_frame:
        text_frame.clear() # 기존 텍스트 모두 삭제
        p = text_frame.paragraphs[0] # 첫 번째 문단 가져오기
        run = p.add_run()
        run.text = new_text

        # 자동 크기 조정 (텍스트가 넘치면 셰이프 크기 조정)
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # 텍스트 정렬 설정
        if align_center:
            p.alignment = PP_ALIGN.CENTER
        else:
            p.alignment = PP_ALIGN.LEFT # 기본은 왼쪽 정렬 (false일 때)

        found_text_frame = True
    
    if not found_text_frame:
        print(f"Warning: Failed to find or edit any suitable text field on slide {slide_index}. No text was inserted.")

    return prs


def _insert_slide_at_index(prs: PresentationType, slide_to_move: SlideType, target_index: int):
    """
    새로 추가된 슬라이드(항상 맨 뒤에 있음)를 원하는 위치에 삽입합니다.
    """
    # 1. 이동할 슬라이드의 내부 XML 요소(sldIdLst 항목)를 찾습니다.
    #    새로 추가된 슬라이드는 항상 _sldIdLst의 가장 마지막에 있습니다.
    xml_slides = prs.slides._sldIdLst
    
    # 슬라이드 ID로 해당하는 xml_slide_element를 찾습니다.
    target_xml_slide_element = None
    for sld_id_elem in xml_slides:
        if sld_id_elem.id == slide_to_move.slide_id:
            target_xml_slide_element = sld_id_elem
            break
    
    if target_xml_slide_element is None:
        print(f"Error: Could not find slide {slide_to_move.slide_id} in internal list for reordering.")
        return

    # 2. 현재 위치에서 해당 XML 요소를 제거합니다.
    xml_slides.remove(target_xml_slide_element)
    
    # 3. 원하는 target_index 위치에 다시 삽입합니다.
    #    prs.slides._sldIdLst는 슬라이드의 순서를 제어하는 내부 리스트입니다.
    xml_slides.insert(target_index, target_xml_slide_element)
    
    print(f"Moved slide {slide_to_move.slide_id} to index {target_index}.")



def add_slides_with_text(
    prs: PresentationType,
    template_slide_index: int,
    texts_for_slides: list,
    is_title_field: bool = False, # 첫 번째 텍스트 필드를 제목으로 사용할지 여부
    target_ph_index: int = None,   # 텍스트를 채울 Placeholder 인덱스 (is_title이 False일 때)
    target_shape_name: str = None # 텍스트를 채울 Shape 이름 (is_title이 False일 때)
) -> dict:
    """
    주어진 텍스트 리스트만큼 슬라이드를 복제하고 텍스트를 채웁니다.
    첫 번째 텍스트는 `template_slide_index`의 원본 슬라이드에 채워지고,
    나머지는 복제된 슬라이드에 채워집니다.
    """
    added_slide_count = 0
    if not texts_for_slides:
        return {"prs": prs, "added_slide_count": 0}

    original_slide_template = prs.slides[template_slide_index]
    slide_layout = original_slide_template.slide_layout

    # 1. 첫 번째 텍스트는 원본 슬라이드에 채웁니다.
    edit_text_field(
        prs=prs,
        slide_index=template_slide_index,
        is_title=is_title_field,
        new_text=texts_for_slides[0],
        ph_index=target_ph_index,
        shape_name=target_shape_name
    )

    # 2. 나머지 텍스트들을 위한 슬라이드를 복제하고 텍스트를 채웁니다.
    for i in range(1, len(texts_for_slides)):
        # 새 슬라이드를 복제하고 콘텐츠 복사 (단순 복제)
        new_slide = prs.slides.add_slide(slide_layout) # 새 슬라이드 추가

        # 새 슬라이드의 위치를 맨 마지막에서 현재 페이지 바로 뒤로 이동
        _insert_slide_at_index(prs, new_slide, template_slide_index + i)

        # 새 슬라이드에 텍스트 채우기
        edit_text_field(
            prs=prs,
            slide_index=prs.slides.index(new_slide), # 새 슬라이드의 실제 인덱스
            is_title=is_title_field,
            new_text=texts_for_slides[i],
            ph_index=target_ph_index,
            shape_name=target_shape_name
        )
        added_slide_count += 1
    
    return {"prs": prs, "added_slide_count": added_slide_count}


# 기존 duplicate_and_add_slide -> add_slides_with_text로 통합/개선
# add_lyrics_slides 함수는 add_slides_with_text의 특정 사용 사례가 됩니다.
def add_lyrics_slides(prs: PresentationType, duplicate_slide_index: int, slide_texts: list) -> dict:
    """
    찬양 가사 슬라이드를 추가합니다.
    `add_slides_with_text`를 사용하여 구현됩니다.
    """
    return add_slides_with_text(prs, duplicate_slide_index, slide_texts, is_title_field=False)


def add_ads_slides(prs: PresentationType, ads_list: list, template_slide_index: int) -> int:
    """
    광고 목록을 기반으로 슬라이드를 추가합니다.
    각 광고는 새 슬라이드에 제목과 내용이 들어갑니다.
    새로 추가되는 슬라이드는 template_slide_index 바로 뒤에 순차적으로 삽입됩니다.
    """
    if not ads_list:
        return 0

    added_count = 0
    original_ad_template_slide = prs.slides[template_slide_index]
    slide_layout = original_ad_template_slide.slide_layout

    for i, ad_data in enumerate(ads_list):
        if i == 0:
            current_slide_index_for_editing = template_slide_index
        else:
            # 새 슬라이드를 추가 (항상 컬렉션 끝에 추가됨)
            new_slide_obj = prs.slides.add_slide(slide_layout)
            
            # 원하는 삽입 위치 계산
            # template_slide_index 바로 뒤에 순차적으로 삽입
            target_insert_index = template_slide_index + i # i는 1부터 시작 (두 번째 슬라이드는 +1, 세 번째는 +2)
            
            # 새 슬라이드를 원하는 위치로 이동
            _insert_slide_at_index(prs, new_slide_obj, target_insert_index)
            
            current_slide_index_for_editing = target_insert_index
            added_count += 1

        # 텍스트 필드 편집
        edit_text_field(
            prs=prs,
            slide_index=current_slide_index_for_editing,
            is_title=True,
            new_text=ad_data.get("title", ""),
            align_center=False
        )
        edit_text_field(
            prs=prs,
            slide_index=current_slide_index_for_editing,
            is_title=False,
            new_text=ad_data.get("contents", ""),
            align_center= False,
            ph_index=1 # 템플릿의 내용 Placeholder 인덱스 (가정)
        )
    return added_count



def add_bible_slides(prs: PresentationType, bible_contents_list: list, template_slide_index: int) -> int:
    """
    성경 구절 내용을 기반으로 슬라이드를 추가합니다.
    각 구절은 새 슬라이드에 제목(구절)과 내용이 들어갑니다.
    새로 추가되는 슬라이드는 template_slide_index 바로 뒤에 순차적으로 삽입됩니다.
    """
    if not bible_contents_list:
        return 0

    added_count = 0
    original_bible_template_slide = prs.slides[template_slide_index]
    slide_layout = original_bible_template_slide.slide_layout

    for i, bible_data in enumerate(bible_contents_list):
        if i == 0:
            current_slide_index_for_editing = template_slide_index
        else:
            # 새 슬라이드를 추가 (항상 컬렉션 끝에 추가됨)
            new_slide_obj = prs.slides.add_slide(slide_layout)
            
            # 원하는 삽입 위치 계산
            target_insert_index = template_slide_index + i
            
            # 새 슬라이드를 원하는 위치로 이동
            _insert_slide_at_index(prs, new_slide_obj, target_insert_index)
            
            current_slide_index_for_editing = target_insert_index
            added_count += 1

        # 텍스트 필드 편집
        edit_text_field(
            prs=prs,
            slide_index=current_slide_index_for_editing,
            is_title=True,
            align_center=False,
            new_text=bible_data.get("title", ""),
        )
        edit_text_field(
            prs=prs,
            slide_index=current_slide_index_for_editing,
            is_title=False,
            new_text=bible_data.get("contents", ""),
            align_center=False,
            ph_index=10 # 템플릿의 내용 Placeholder 인덱스 (가정)
        )
    return added_count


def save_presentation(prs: PresentationType, save_path: str):
    """
    프레젠테이션을 지정된 경로에 저장합니다.
    """
    try:
        prs.save(save_path)
        print(f"Presentation saved to: {save_path}")
    except Exception as e:
        print(f"Error saving presentation to {save_path}: {e}")
        raise